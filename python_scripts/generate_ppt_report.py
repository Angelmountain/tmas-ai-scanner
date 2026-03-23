#!/usr/bin/env python3
"""
Generate PowerPoint report from security assessment Excel data.

Reads Excel files produced by the security assessment pipeline and updates
chart data in a PowerPoint template, creating a branded Trend Micro
NDR Security Assessment report.

The template uses 3D bar charts (bar3DChart).  python-pptx cannot read
their chart_type property, but *can* replace the underlying data via
``chart.replace_data()``.  This script therefore NEVER deletes or recreates
charts -- it only pushes new CategoryChartData into the existing chart
objects, preserving all styling, colours, 3D effects, and layout.

Usage (CLI):
    python generate_ppt_report.py \\
        --template templates/NDR_Security_Assessment_v2.0.pptx \\
        --data-dir output_directory \\
        --output reports/final_report.pptx

Usage (module):
    from generate_ppt_report import generate_report
    path = generate_report(
        template_path="templates/NDR_Security_Assessment_v2.0.pptx",
        data_dir="output_directory",
        output_path="reports/final_report.pptx",
    )
"""

import argparse
import logging
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import openpyxl
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Trend Micro brand colours
TREND_COLORS = {
    "red": RGBColor(0xD5, 0x2B, 0x1E),       # #D52B1E
    "dark_gray": RGBColor(0x4A, 0x4A, 0x4A),  # #4A4A4A
    "light_gray": RGBColor(0x9B, 0x9B, 0x9B),
    "blue": RGBColor(0x00, 0x7C, 0xB0),       # #007CB0
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

# ---------------------------------------------------------------------------
# Slide mapping for the v2.0 template (44 slides)
#
# Keys   = search / topic names that match Excel file stems.
# Values = 1-indexed slide number in the template, or None when the
#          search has no dedicated slide (data lives in Excel only).
#
# The template slide titles do not always match the search names
# exactly (e.g. slide 24 is labelled "PUA country connections" but
# maps to the "PUA VPN Services" search).  The mapping below is the
# authoritative source of truth for the v2.0 template.
# ---------------------------------------------------------------------------

SLIDE_MAPPING: Dict[str, Optional[int]] = {
    "Network Detections": 6,
    "Top Accounts used": 7,
    "Server ports used": 8,
    "Unsuccessful logon": 9,
    "Top File used": None,           # No slide in v2.0 template
    "Top File types used": None,     # No slide in v2.0 template
    "Protocols used": 12,
    "Request methods": 13,
    "Response codes": 14,
    "SSL Cert Common Name": 15,
    "SSH Detections": 17,
    "SSH Versions": 18,
    "PUA AI Services": 20,
    "PUA root detections": 21,
    "PUA Remote Access": 22,
    "PUA Cloud Storage": 23,
    "PUA VPN Services": 24,          # slide title: "PUA country connections"
    "PUA Pastebin": 25,              # slide title: "PUA email attachments"
    "PUA Darknet links": 26,
    "PUA Administrator Usage": 27,
    "RDP User Usage": 29,
    "RDP Source IP": 30,
    "RDP Destination IP": 31,
    "Suspicious TLDs": 33,           # slide title: ".CH requests"
    "Bad States": 34,                # slide title: ".RU requests"
    "Russian IT-companies": 35,
    "EPP/EDR/XDR Vendors": 37,
    "Firewall Vendors": 38,
    "US Vendors": 39,
}

# Maximum categories shown per chart / fallback table
MAX_CHART_CATEGORIES = 15
MAX_TABLE_ROWS = 10

# Default template location (relative to repository root)
DEFAULT_TEMPLATE = "templates/NDR_Security_Assessment_v2.0.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _normalise(name: str) -> str:
    """Lowercase, strip whitespace, remove non-alphanumerics."""
    return re.sub(r"[^a-z0-9]", "", name.lower())


# ---------------------------------------------------------------------------
# PowerPointReportGenerator
# ---------------------------------------------------------------------------

class PowerPointReportGenerator:
    """Load a .pptx template, inject chart data from Excel files, and save."""

    def __init__(self, template_path: Path, output_path: Path):
        self.template_path = template_path
        self.output_path = output_path
        self.prs: Optional[Presentation] = None
        self._stats: Dict[str, str] = {}

    # -- lifecycle ----------------------------------------------------------

    def load_template(self) -> None:
        """Open the PowerPoint template."""
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")
        logger.info("Loading template: %s", self.template_path)
        self.prs = Presentation(str(self.template_path))
        logger.info("Template loaded with %d slides", len(self.prs.slides))

    def save(self) -> Path:
        """Write the updated presentation to *output_path*."""
        if self.prs is None:
            raise RuntimeError("No presentation loaded; call load_template() first")
        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(self.output_path))
        logger.info("Report saved to: %s", self.output_path)
        return self.output_path

    # -- data ingestion -----------------------------------------------------

    @staticmethod
    def read_excel_data(excel_path: Path) -> List[Tuple[str, int]]:
        """Return up to *MAX_CHART_CATEGORIES* (category, count) pairs.

        Expects two columns: Category (col A) and Count (col B), with a
        header row.  Rows with missing or non-positive counts are skipped.
        """
        try:
            wb = openpyxl.load_workbook(excel_path, data_only=True)
            ws = wb.active

            data: List[Tuple[str, int]] = []
            for row in ws.iter_rows(min_row=2, values_only=True):
                if row is None or len(row) < 2:
                    continue
                category = str(row[0]).strip() if row[0] else "Unknown"
                try:
                    value = int(row[1]) if row[1] is not None else 0
                except (ValueError, TypeError):
                    value = 0
                if category and value > 0:
                    data.append((category, value))

            wb.close()
            return data[:MAX_CHART_CATEGORIES]

        except Exception:
            logger.exception("Failed to read Excel file %s", excel_path)
            return []

    @staticmethod
    def find_excel_file(search_name: str, data_dir: Path) -> Optional[Path]:
        """Locate the Excel file for *search_name* inside *data_dir*.

        Tries several strategies in order:

        1. Exact prefix glob with common suffixes (_horizontal, _vertical).
        2. Variant with spaces replaced by underscores (and vice-versa).
        3. Normalised substring match (strip all non-alphanum, compare).
        """
        if not data_dir.is_dir():
            return None

        # Build candidate prefixes: original name, underscore variant, space variant
        variants = {search_name}
        variants.add(search_name.replace(" ", "_"))
        variants.add(search_name.replace("_", " "))
        # Also try with slashes removed (some filenames drop punctuation)
        variants.add(search_name.replace("/", "").replace("\\", ""))

        suffixes = ["*_horizontal.xlsx", "*_vertical.xlsx", "*.xlsx"]

        for variant in variants:
            for suffix in suffixes:
                pattern = f"{variant}{suffix}"
                matches = sorted(data_dir.glob(pattern))
                if matches:
                    return matches[0]

        # Case-insensitive exact-prefix: retry with lowered glob isn't
        # possible on all OS, so do a manual scan instead.
        search_lower = search_name.lower()
        search_under = search_lower.replace(" ", "_")
        search_space = search_lower.replace("_", " ")

        all_xlsx = sorted(data_dir.glob("*.xlsx"))

        for xlsx_file in all_xlsx:
            stem_lower = xlsx_file.stem.lower()
            if (
                stem_lower.startswith(search_lower)
                or stem_lower.startswith(search_under)
                or stem_lower.startswith(search_space)
            ):
                return xlsx_file

        # Fuzzy fallback: normalised substring match
        clean = _normalise(search_name)
        for xlsx_file in all_xlsx:
            stem_clean = _normalise(xlsx_file.stem)
            if clean in stem_clean or stem_clean in clean:
                return xlsx_file

        return None

    # -- slide manipulation -------------------------------------------------

    def update_chart_on_slide(
        self, slide_num: int, data: List[Tuple[str, int]], title: str
    ) -> bool:
        """Replace chart data on slide *slide_num* (1-indexed).

        The template uses 3D bar charts (bar3DChart).  python-pptx cannot
        read their ``chart_type`` but CAN replace their data through the
        ``chart.replace_data()`` API which writes directly to the embedded
        Excel workbook.  This preserves all chart formatting, colours, 3D
        effects, and animations.

        If ``replace_data`` raises (e.g. a genuinely incompatible chart
        type), a small data table is added below the chart as a fallback.
        The chart itself is NEVER deleted or recreated.

        Returns True when the slide was successfully updated (either via
        chart data replacement or fallback table).
        """
        if self.prs is None:
            raise RuntimeError("No presentation loaded")

        total_slides = len(self.prs.slides)
        if slide_num < 1 or slide_num > total_slides:
            logger.warning(
                "Slide %d out of range (template has %d slides) for '%s'",
                slide_num,
                total_slides,
                title,
            )
            return False

        slide = self.prs.slides[slide_num - 1]

        # Locate the first chart shape on the slide
        chart_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "has_chart") and shape.has_chart:
                chart_shape = shape
                break

        categories = [item[0] for item in data]
        values = [item[1] for item in data]

        if chart_shape is not None:
            try:
                chart = chart_shape.chart
                chart_data = CategoryChartData()
                chart_data.categories = categories
                chart_data.add_series("Count", values)
                chart.replace_data(chart_data)
                logger.info(
                    "Replaced chart data on slide %d for '%s' (%d categories)",
                    slide_num,
                    title,
                    len(data),
                )
                return True
            except Exception:
                logger.warning(
                    "chart.replace_data() failed on slide %d for '%s'; "
                    "falling back to data table",
                    slide_num,
                    title,
                    exc_info=True,
                )
                # Do NOT delete or recreate the chart -- just add a table below
                self._add_data_table(slide, data, title)
                return True

        # No chart shape found on the slide -- add a table instead
        logger.info(
            "No chart found on slide %d for '%s'; adding data table",
            slide_num,
            title,
        )
        self._add_data_table(slide, data, title)
        return True

    # -- data table fallback ------------------------------------------------

    @staticmethod
    def _add_data_table(
        slide, data: List[Tuple[str, int]], title: str
    ) -> None:
        """Insert a small branded data table near the bottom of *slide*."""
        data = data[:MAX_TABLE_ROWS]
        if not data:
            logger.warning("No data to add table for '%s'", title)
            return

        rows = len(data) + 1  # +1 for header
        cols = 2
        left = Inches(0.5)
        top = Inches(4.5)
        width = Inches(5.0)
        height = Inches(0.3 * rows)

        # Clamp height so the table does not overflow the slide
        max_height = Inches(2.5)
        if height > max_height:
            height = max_height

        try:
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            # Header row
            table.cell(0, 0).text = "Category"
            table.cell(0, 1).text = "Count"
            for col_idx in range(cols):
                cell = table.cell(0, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = TREND_COLORS["red"]
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.color.rgb = TREND_COLORS["white"]
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(10)

            # Data rows
            for row_idx, (category, value) in enumerate(data, start=1):
                cat_text = (
                    str(category)[:40] + "..."
                    if len(str(category)) > 40
                    else str(category)
                )
                table.cell(row_idx, 0).text = cat_text
                table.cell(row_idx, 1).text = str(value)

                for col_idx in range(cols):
                    cell = table.cell(row_idx, col_idx)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(9)
                        paragraph.font.color.rgb = TREND_COLORS["dark_gray"]

            # Column widths
            table.columns[0].width = Inches(3.5)
            table.columns[1].width = Inches(1.5)

            logger.info(
                "Added data table (%d rows) to slide for '%s'",
                len(data),
                title,
            )
        except Exception:
            logger.exception("Failed to add data table for '%s'", title)

    # -- orchestration ------------------------------------------------------

    def process_all(self, data_dir: Path) -> Dict[str, str]:
        """Walk *SLIDE_MAPPING* and update every slide that has matching data.

        The *data_dir* may contain a ``ppt/`` subdirectory; if present it
        is searched first.

        Returns a dict mapping search names to a status string:
        ``"chart_updated"``, ``"table_added"``, ``"no_data"``,
        ``"no_file"``, ``"skipped"``, ``"out_of_range"``.
        """
        results: Dict[str, str] = {}

        # Prefer a "ppt" sub-folder if it exists
        ppt_dir = data_dir / "ppt"
        if not ppt_dir.is_dir():
            ppt_dir = data_dir

        logger.info("Searching for Excel files in: %s", ppt_dir)

        for search_name, slide_num in SLIDE_MAPPING.items():
            # Locate Excel file (try ppt_dir first, then data_dir)
            excel_file = self.find_excel_file(search_name, ppt_dir)
            if excel_file is None and ppt_dir != data_dir:
                excel_file = self.find_excel_file(search_name, data_dir)

            if excel_file is None:
                logger.debug("No Excel file found for '%s'", search_name)
                results[search_name] = "no_file"
                continue

            logger.info("Found '%s' -> %s", search_name, excel_file.name)

            data = self.read_excel_data(excel_file)
            if not data:
                logger.warning(
                    "Excel file for '%s' contains no usable data", search_name
                )
                results[search_name] = "no_data"
                continue

            if slide_num is None:
                logger.info(
                    "'%s' has no slide mapping; data available in Excel only",
                    search_name,
                )
                results[search_name] = "skipped"
                continue

            success = self.update_chart_on_slide(slide_num, data, search_name)
            results[search_name] = "updated" if success else "out_of_range"

        self._stats = results
        return results

    def print_summary(self) -> None:
        """Log a human-readable summary of processing results."""
        if not self._stats:
            return

        updated = sum(1 for v in self._stats.values() if v == "updated")
        no_file = sum(1 for v in self._stats.values() if v == "no_file")
        no_data = sum(1 for v in self._stats.values() if v == "no_data")
        skipped = sum(1 for v in self._stats.values() if v == "skipped")
        failed = sum(
            1 for v in self._stats.values() if v == "out_of_range"
        )

        total_mapped = sum(1 for v in SLIDE_MAPPING.values() if v is not None)

        logger.info("=" * 60)
        logger.info("Report generation summary")
        logger.info("-" * 60)
        logger.info("  Slides updated:       %d / %d", updated, total_mapped)
        logger.info("  No Excel file found:  %d", no_file)
        logger.info("  Excel file empty:     %d", no_data)
        logger.info("  No slide (skipped):   %d", skipped)
        if failed:
            logger.warning("  Out of range/failed:  %d", failed)
        logger.info("=" * 60)


# ---------------------------------------------------------------------------
# Public API for module-level imports
# ---------------------------------------------------------------------------

def generate_report(
    template_path: str,
    data_dir: str,
    output_path: Optional[str] = None,
    *,
    log_level: str = "INFO",
) -> Optional[Path]:
    """High-level helper: generate a PowerPoint report in one call.

    Parameters
    ----------
    template_path:
        Path to the ``.pptx`` template (v2.0 with 44 slides).
    data_dir:
        Directory containing Excel files (may have a ``ppt/`` sub-folder).
    output_path:
        Destination ``.pptx`` file.  If *None* a timestamped file is
        written next to the template.
    log_level:
        Python log-level name (``DEBUG``, ``INFO``, ``WARNING``, ``ERROR``).

    Returns
    -------
    pathlib.Path to the written report, or *None* on failure.
    """
    _configure_logging(log_level)

    template = Path(template_path)
    data = Path(data_dir)

    if output_path is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output = template.parent / f"NDR_Security_Assessment_Report_{timestamp}.pptx"
    else:
        output = Path(output_path)

    try:
        gen = PowerPointReportGenerator(template, output)
        gen.load_template()
        gen.process_all(data)
        gen.print_summary()
        return gen.save()

    except Exception:
        logger.exception("Report generation failed")
        return None


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _configure_logging(level_name: str = "INFO") -> None:
    """Set up root logger with a consistent format."""
    level = getattr(logging, level_name.upper(), logging.INFO)
    logging.basicConfig(
        level=level,
        format="%(asctime)s  %(name)s  %(levelname)-8s  %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S",
    )


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Generate a Trend Micro NDR Security Assessment PowerPoint report "
            "from Excel data files.  Uses the v2.0 template (44 slides) with "
            "3D bar charts -- chart data is replaced in-place, preserving all "
            "styling and 3D effects."
        ),
    )
    parser.add_argument(
        "--template",
        default=DEFAULT_TEMPLATE,
        help=(
            "Path to the .pptx template "
            f"(default: {DEFAULT_TEMPLATE})"
        ),
    )
    parser.add_argument(
        "--data-dir",
        required=True,
        help="Directory containing the Excel (.xlsx) files with assessment data.",
    )
    parser.add_argument(
        "--output",
        default=None,
        help=(
            "Output .pptx file path.  Defaults to a timestamped file in "
            "the same directory as the template."
        ),
    )
    parser.add_argument(
        "--log-level",
        default="INFO",
        choices=["DEBUG", "INFO", "WARNING", "ERROR"],
        help="Logging verbosity (default: INFO).",
    )
    return parser


def main(argv: Optional[List[str]] = None) -> int:
    """CLI entry point.  Returns 0 on success, 1 on failure."""
    parser = _build_parser()
    args = parser.parse_args(argv)

    result = generate_report(
        template_path=args.template,
        data_dir=args.data_dir,
        output_path=args.output,
        log_level=args.log_level,
    )

    if result is None:
        logger.error("Report generation failed.")
        return 1

    print(f"Report written to: {result}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
